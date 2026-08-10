"""Builds a persuasive, diagram-rich academic pitch deck."""
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---------- Design system (matches the product UI palette) ----------
BG1      = RGBColor(0x02, 0x15, 0x30)   # navy-deep (app background)
BG2      = RGBColor(0x0a, 0x2d, 0x5e)   # navy-mid
PANEL    = RGBColor(0x0a, 0x26, 0x48)   # s2 surface
PANEL2   = RGBColor(0x0f, 0x32, 0x5a)   # s3 surface
PRIMARY  = RGBColor(0x03, 0x1b, 0x47)   # navy
ACCENT   = RGBColor(0xee, 0x89, 0x23)   # brand orange
ACCENT2  = RGBColor(0x5b, 0x9b, 0xd5)   # brand blue
OKGREEN  = RGBColor(0x5c, 0xb8, 0x8a)   # ok
ERRRED   = RGBColor(0xe8, 0x55, 0x55)   # err
GOLD     = RGBColor(0xe8, 0xa8, 0x38)   # warn / third accent
WHITE    = RGBColor(0xe8, 0xee, 0xf7)   # txt
MUTED    = RGBColor(0x8a, 0xa3, 0xc4)   # muted
SOFT     = RGBColor(0xc8, 0xd6, 0xe5)   # soft light
INK      = RGBColor(0x02, 0x15, 0x30)   # dark text on accent fills

FONT = "Segoe UI"
TOTAL = 17

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def notes(s, body):
    s.notes_slide.notes_text_frame.text = body


def _to_back(s, sp):
    el = sp._element
    el.getparent().remove(el)
    s.shapes._spTree.insert(2, el)


def set_bg(s, c1=BG1, c2=BG2, angle=70):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.line.fill.background()
    r.shadow.inherit = False
    f = r.fill
    f.gradient()
    f.gradient_angle = angle
    f.gradient_stops[0].color.rgb = c1
    f.gradient_stops[0].position = 0.0
    f.gradient_stops[1].color.rgb = c2
    f.gradient_stops[1].position = 1.0
    _to_back(s, r)
    return r


def _soft_shadow(sp):
    spPr = sp._element.spPr
    eff = spPr.find(qn('a:effectLst'))
    if eff is None:
        eff = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(eff, qn('a:outerShdw'))
    sh.set('blurRad', '180000'); sh.set('dist', '70000')
    sh.set('dir', '5400000'); sh.set('rotWithShape', '0')
    c = etree.SubElement(sh, qn('a:srgbClr')); c.set('val', '000000')
    a = etree.SubElement(c, qn('a:alpha')); a.set('val', '42000')


def _set_alpha(sp, opacity_pct):
    srgb = sp.fill.fore_color._xFill.find(qn('a:srgbClr'))
    if srgb is not None:
        a = etree.SubElement(srgb, qn('a:alpha'))
        a.set('val', str(int(opacity_pct * 1000)))


def _dash(sp):
    ln = sp.line._get_or_add_ln()
    d = etree.SubElement(ln, qn('a:prstDash')); d.set('val', 'dash')


def rect(s, x, y, w, h, color=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE,
         shadow=False, alpha=None, dash=False):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w or Pt(1)
    if shadow:
        _soft_shadow(sp)
    if alpha is not None and color is not None:
        _set_alpha(sp, alpha)
    if dash and line is not None:
        _dash(sp)
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, italic, font) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
    return tb


def R(txt, size, color=WHITE, bold=False, italic=False, font=FONT):
    return (txt, size, color, bold, italic, font)


def connect(s, x1, y1, x2, y2, color=MUTED, width=Pt(1.75), arrow=True, dash=False):
    cx = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cx.shadow.inherit = False
    cx.line.color.rgb = color
    cx.line.width = width
    ln = cx.line._get_or_add_ln()
    if dash:
        d = etree.SubElement(ln, qn('a:prstDash')); d.set('val', 'dash')
    if arrow:
        te = etree.SubElement(ln, qn('a:tailEnd'))
        te.set('type', 'triangle'); te.set('w', 'med'); te.set('len', 'med')
    return cx


def chevron(s, x, y, w=Inches(0.22), h=Inches(0.4), color=MUTED):
    a = rect(s, x, y, w, h, shape=MSO_SHAPE.CHEVRON)
    a.fill.background()
    a.line.color.rgb = color
    a.line.width = Pt(2.25)
    return a


def add_transition(s, kind="morph", speed="slow", dur_ms=700):
    sld = s._element
    _P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    _MC = "{http://schemas.openxmlformats.org/markup-compatibility/2006}"
    for tag in (_P + "transition", _MC + "AlternateContent"):
        for el in list(sld):
            if el.tag == tag:
                sld.remove(el)
    cSld = sld.find(qn('p:cSld'))
    if kind == "morph":
        xml = (
            '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">'
            '<mc:Choice xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" Requires="p14">'
            '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'spd="%s" p14:dur="%d"><p14:morph option="byObject"/></p:transition></mc:Choice>'
            '<mc:Fallback><p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'spd="%s"><p:fade/></p:transition></mc:Fallback></mc:AlternateContent>'
            % (speed, dur_ms, speed)
        )
    else:
        inner = {"fade": '<p:fade/>', "push": '<p:push dir="l"/>'}.get(kind, '<p:fade/>')
        xml = (
            '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">'
            '<mc:Choice xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" Requires="p14">'
            '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'spd="%s" p14:dur="%d">%s</p:transition></mc:Choice>'
            '<mc:Fallback><p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'spd="%s">%s</p:transition></mc:Fallback></mc:AlternateContent>'
            % (speed, dur_ms, inner, speed, inner)
        )
    cSld.addnext(etree.fromstring(xml))


def kicker(s, label, num):
    rect(s, Inches(0.6), Inches(0.5), Inches(0.07), Inches(0.42), color=ACCENT)
    text(s, Inches(0.8), Inches(0.45), Inches(9), Inches(0.5),
         [[R(label.upper(), 13, ACCENT, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, Inches(11.9), Inches(0.45), Inches(0.9), Inches(0.5),
         [[R("%02d" % num, 12, MUTED), R(" / %d" % TOTAL, 12, MUTED)]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def title(s, lines, y=Inches(1.0), size=29, w=Inches(11.9), color=WHITE):
    text(s, Inches(0.8), y, w, Inches(1.4),
         [[R(ln, size, color, True)] for ln in lines], line_spacing=1.05, space_after=2)


def node(s, x, y, w, h, head, body=None, color=ACCENT, head_size=15, body_size=11.5,
         fill=PANEL, head_color=WHITE):
    rect(s, x, y, w, h, color=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, y, Inches(0.08), h, color=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if body:
        text(s, x + Inches(0.22), y + Inches(0.14), w - Inches(0.4), Inches(0.45),
             [[R(head, head_size, head_color, True)]], space_after=0)
        text(s, x + Inches(0.22), y + Inches(0.14) + Inches(0.42), w - Inches(0.42), h - Inches(0.7),
             [[R(body, body_size, MUTED)]], line_spacing=1.05, space_after=0)
    else:
        text(s, x + Inches(0.22), y, w - Inches(0.4), h,
             [[R(head, head_size, head_color, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def chip(s, x, y, w, label, color=ACCENT2, h=Inches(0.4), size=11.5, txt=None):
    rect(s, x, y, w, h, color=PANEL2, line=color, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, y, w, h, [[R(label, size, txt or SOFT, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def img_placeholder(s, x, y, w, h, caption, art="Illustration"):
    rect(s, x, y, w, h, color=PANEL, line=ACCENT2, line_w=Pt(1.25),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, dash=True)
    # framed-picture glyph
    rect(s, x + w / 2 - Inches(0.45), y + Inches(0.38), Inches(0.9), Inches(0.62),
         color=None, line=ACCENT2, line_w=Pt(1.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x + w / 2 - Inches(0.27), y + Inches(0.74), Inches(0.22), Inches(0.22),
         color=ACCENT2, shape=MSO_SHAPE.OVAL)
    text(s, x + Inches(0.2), y + h / 2 - Inches(0.05), w - Inches(0.4), Inches(0.45),
         [[R(art.upper() + " PLACEHOLDER", 12, ACCENT2, True)]],
         align=PP_ALIGN.CENTER, space_after=0)
    text(s, x + Inches(0.3), y + h - Inches(1.0), w - Inches(0.6), Inches(0.85),
         [[R(caption, 11, MUTED, italic=True)]], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.BOTTOM, line_spacing=1.1, space_after=0)


def big_stat(s, x, y, w, value, label, color=ACCENT, h=Inches(2.0), vsize=44):
    rect(s, x, y, w, h, color=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, y, w, Inches(0.1), color=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, y + Inches(0.34), w, Inches(0.95),
         [[R(value, vsize, color, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x + Inches(0.18), y + h - Inches(0.78), w - Inches(0.36), Inches(0.7),
         [[R(label, 12.5, SOFT)]], align=PP_ALIGN.CENTER, line_spacing=1.05, space_after=0)


# =================================================================
# 1 - TITLE (hero)
# =================================================================
s = slide()
set_bg(s, BG1, BG2, angle=60)
rect(s, Inches(9.3), Inches(-1.7), Inches(6.2), Inches(6.2), color=ACCENT, alpha=8, shape=MSO_SHAPE.OVAL)
rect(s, Inches(10.8), Inches(3.7), Inches(5.2), Inches(5.2), color=ACCENT2, alpha=7, shape=MSO_SHAPE.OVAL)
text(s, Inches(0.9), Inches(1.2), Inches(11), Inches(0.6),
     [[R("INFRASTRUCTURE AS CODE  x  ARTIFICIAL INTELLIGENCE", 14, ACCENT, True)]], space_after=0)
rect(s, Inches(0.92), Inches(1.78), Inches(0.9), Inches(0.10), color=ACCENT)
text(s, Inches(0.85), Inches(2.05), Inches(8.8), Inches(2.6),
     [[R("Generic AI guesses", 46, WHITE, True)],
      [R("your infrastructure.", 46, WHITE, True)],
      [R("We prove it.", 46, ACCENT, True)]],
     line_spacing=1.02, space_after=2)
text(s, Inches(0.9), Inches(4.95), Inches(8.7), Inches(1.1),
     [[R("An AI assistant that writes production-ready automation from plain English - "
         "grounded in official documentation and verified before you ever run it.", 16, SOFT)]],
     line_spacing=1.25, space_after=0)
rect(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.02), color=PANEL2)
text(s, Inches(0.9), Inches(6.42), Inches(8), Inches(0.5),
     [[R("Final-Year Engineering Project (PFE)  -  2026", 13, MUTED)]], space_after=0)
text(s, Inches(9.4), Inches(6.42), Inches(3.0), Inches(0.5),
     [[R("Project Defense", 13, MUTED)]], align=PP_ALIGN.RIGHT, space_after=0)
add_transition(s, "fade")

# =================================================================
# 2 - STAKES / HOOK
# =================================================================
s = slide()
set_bg(s)
kicker(s, "Why this matters", 2)
title(s, ["Infrastructure as Code runs the modern cloud -",
          "and a single wrong line can bring it down"])
text(s, Inches(0.8), Inches(2.2), Inches(7.2), Inches(0.9),
     [[R("Today every deployment, server and cloud resource is described as code. "
         "That code is powerful - and unforgiving. The cost of a mistake is not a typo; "
         "it is downtime, a security hole, or hours of lost engineering time.", 15, SOFT)]],
     line_spacing=1.25, space_after=0)
stakes = [
    ("Wrong module", "A plausible-looking task that simply does not exist - the deploy fails.", ERRRED),
    ("Wrong parameter", "A missing or misused option silently opens a security gap.", GOLD),
    ("Lost hours", "Engineers burn time debugging YAML instead of shipping.", ACCENT2),
]
for i, (h, b, c) in enumerate(stakes):
    y = Inches(3.45) + i * Inches(1.12)
    rect(s, Inches(0.8), y, Inches(7.2), Inches(0.95), color=PANEL,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, Inches(0.8), y, Inches(0.09), Inches(0.95), color=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(1.1), y, Inches(2.4), Inches(0.95),
         [[R(h, 16, c, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, Inches(3.4), y, Inches(4.45), Inches(0.95),
         [[R(b, 12.5, SOFT)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1, space_after=0)
img_placeholder(s, Inches(8.45), Inches(2.2), Inches(4.05), Inches(4.35),
                "Suggested: a screenshot of a failed CI/CD pipeline or a red Ansible "
                "error trace - something that visually conveys \"broken deploy\".",
                art="Screenshot")
add_transition(s, "morph")

# =================================================================
# 3 - THE PROBLEMATIC
# =================================================================
s = slide()
set_bg(s)
kicker(s, "The Problematic", 3)
title(s, ["The research question"])
rect(s, Inches(1.1), Inches(2.15), Inches(11.1), Inches(2.0), color=PANEL2,
     line=ACCENT, line_w=Pt(1.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
text(s, Inches(1.5), Inches(2.35), Inches(10.3), Inches(1.6),
     [[R("How can we let ", 21, WHITE), R("anyone", 21, ACCENT, True),
       R(" generate ", 21, WHITE), R("correct, secure, production-grade", 21, ACCENT, True),
       R(" Infrastructure as Code", 21, WHITE)],
      [R("from plain language - ", 21, WHITE),
       R("without sacrificing reliability, trust, or expertise?", 21, ACCENT, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2, space_after=4)
subs = [
    ("Correctness", "Will the generated code actually run, with valid modules and parameters?"),
    ("Trust", "Can we explain and audit where every answer came from?"),
    ("Accessibility", "Can a non-expert get a safe result without years of experience?"),
]
for i, (h, b) in enumerate(subs):
    x = Inches(0.8) + i * Inches(4.05)
    rect(s, x, Inches(4.65), Inches(3.8), Inches(1.95), color=PANEL,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, Inches(4.65), Inches(3.8), Inches(0.1), color=ACCENT2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + Inches(0.3), Inches(4.9), Inches(3.2), Inches(0.5),
         [[R(h, 17, ACCENT2, True)]], space_after=0)
    text(s, x + Inches(0.3), Inches(5.45), Inches(3.25), Inches(1.1),
         [[R(b, 13, SOFT)]], line_spacing=1.15, space_after=0)
add_transition(s, "morph")

# =================================================================
# 4 - WHY NOT CHATGPT / CLAUDE
# =================================================================
s = slide()
set_bg(s)
kicker(s, "The Gap in Today's Tools", 4)
title(s, ["General chatbots can write YAML -",
          "but they cannot be trusted to run it"])
text(s, Inches(0.8), Inches(2.15), Inches(11.6), Inches(0.55),
     [[R("Tools like ChatGPT and Claude answer from memory. For high-stakes, "
         "fast-moving infrastructure, that memory is a liability:", 15, SOFT)]],
     line_spacing=1.2, space_after=0)
gaps = [
    ("Hallucinated modules", "Invents tasks and options that look right but do not exist."),
    ("No grounding", "Not connected to the current, official documentation."),
    ("Knowledge cutoff", "Frozen at training time; unaware of new or changed modules."),
    ("No validation", "You only discover the errors at runtime, in production."),
    ("No provenance", "Cannot show or justify where an answer came from."),
    ("Not specialized", "A generalist guessing in a domain that punishes guesses."),
]
cw = Inches(3.78); ch = Inches(1.4); gx = Inches(0.25); gy = Inches(0.22)
for i, (h, b) in enumerate(gaps):
    col = i % 3; row = i // 3
    x = Inches(0.8) + col * (cw + gx)
    y = Inches(2.95) + row * (ch + gy)
    rect(s, x, y, cw, ch, color=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x + Inches(0.22), y + Inches(0.24), Inches(0.4), Inches(0.4),
         color=ERRRED, shape=MSO_SHAPE.OVAL)
    text(s, x + Inches(0.22), y + Inches(0.22), Inches(0.4), Inches(0.4),
         [[R("X", 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x + Inches(0.78), y + Inches(0.2), cw - Inches(0.95), Inches(0.45),
         [[R(h, 14.5, WHITE, True)]], space_after=0)
    text(s, x + Inches(0.78), y + Inches(0.66), cw - Inches(0.95), Inches(0.65),
         [[R(b, 11.5, MUTED)]], line_spacing=1.05, space_after=0)
text(s, Inches(0.8), Inches(6.55), Inches(11.6), Inches(0.5),
     [[R("The problem is not the model's fluency - it is the ", 14, SOFT),
       R("absence of evidence and verification.", 14, ACCENT, True)]], space_after=0)
add_transition(s, "morph")

# =================================================================
# 5 - VALUE PROPOSITION (3 pillars)
# =================================================================
s = slide()
set_bg(s)
kicker(s, "Our Answer", 5)
title(s, ["We turn a language model into a domain expert you can trust"])
pillars = [
    ("GROUNDED", "It retrieves the exact official documentation before writing a single line - "
     "so answers are built on evidence, not memory.", ACCENT),
    ("VERIFIED", "Every playbook passes a multi-layer validation gate - so errors are caught "
     "before they ever reach you.", OKGREEN),
    ("TRANSPARENT", "It cites the precise modules behind each answer - so every result is "
     "explainable and auditable.", ACCENT2),
]
cw = Inches(3.8); gap = Inches(0.25)
for i, (h, b, c) in enumerate(pillars):
    x = Inches(0.8) + i * (cw + gap)
    y = Inches(2.45)
    rect(s, x, y, cw, Inches(3.9), color=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, y, cw, Inches(0.85), color=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y + Inches(0.45), cw, Inches(0.4), color=c)
    text(s, x, y + Inches(0.16), cw, Inches(0.55),
         [[R(h, 20, INK, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x + Inches(0.35), y + Inches(1.15), cw - Inches(0.7), Inches(2.5),
         [[R(b, 14.5, SOFT)]], line_spacing=1.3, space_after=0)
text(s, Inches(0.8), Inches(6.6), Inches(11.6), Inches(0.5),
     [[R("Grounded. Verified. Transparent. ", 15, ACCENT, True),
       R("The three things a general chatbot cannot promise.", 15, SOFT)]],
     align=PP_ALIGN.CENTER, space_after=0)
add_transition(s, "morph")

# =================================================================
# 6 - THE CORE DIFFERENCE (diagram: generic vs grounded)
# =================================================================
s = slide()
set_bg(s)
kicker(s, "How We Are Different", 6)
title(s, ["Same question, two very different paths"])
# Generic path (top)
text(s, Inches(0.8), Inches(2.0), Inches(5), Inches(0.4),
     [[R("A GENERAL CHATBOT", 13, ERRRED, True)]], space_after=0)
gy1 = Inches(2.5)
node(s, Inches(0.8), gy1, Inches(2.5), Inches(0.95), "Your request", color=MUTED, head_size=13)
node(s, Inches(4.4), gy1, Inches(3.0), Inches(0.95), "LLM from memory",
     "answers from training data", color=ERRRED, head_size=14)
node(s, Inches(8.7), gy1, Inches(3.8), Inches(0.95), "Plausible answer",
     "may be wrong - found out at runtime", color=ERRRED, head_size=14)
connect(s, Inches(3.3), gy1 + Inches(0.47), Inches(4.4), gy1 + Inches(0.47), color=MUTED)
connect(s, Inches(7.4), gy1 + Inches(0.47), Inches(8.7), gy1 + Inches(0.47), color=ERRRED)
# divider
rect(s, Inches(0.8), Inches(3.95), Inches(11.7), Inches(0.015), color=PANEL2)
# Grounded path (bottom)
text(s, Inches(0.8), Inches(4.15), Inches(6), Inches(0.4),
     [[R("OUR GROUNDED, VERIFIED APPROACH", 13, OKGREEN, True)]], space_after=0)
gy2 = Inches(4.65)
steps = [
    ("Your request", None, MUTED, Inches(1.95)),
    ("Retrieve docs", "official sources", ACCENT, Inches(2.05)),
    ("Agent generates", "from evidence", ACCENT, Inches(2.05)),
    ("Validate", "catch errors", OKGREEN, Inches(1.85)),
    ("Trusted result", "+ cited sources", OKGREEN, Inches(2.15)),
]
x = Inches(0.8)
positions = []
for (h, b, c, w) in steps:
    node(s, x, gy2, w, Inches(1.05), h, b, color=c, head_size=13.5, body_size=11)
    positions.append((x, w))
    x = x + w + Inches(0.22)
for i in range(len(positions) - 1):
    x0, w0 = positions[i]
    x1, _ = positions[i + 1]
    col = OKGREEN if i >= 2 else ACCENT
    connect(s, x0 + w0, gy2 + Inches(0.52), x1, gy2 + Inches(0.52), color=col)
text(s, Inches(0.8), Inches(6.5), Inches(11.6), Inches(0.5),
     [[R("Evidence in, verification out. ", 14, OKGREEN, True),
       R("The detour through documentation and validation is the whole point.", 14, SOFT)]],
     space_after=0)
add_transition(s, "morph")

# =================================================================
# 7 - ARCHITECTURE (diagram)
# =================================================================
s = slide()
set_bg(s)
kicker(s, "Architecture", 7)
title(s, ["Three cooperating layers turn a question into a trustworthy answer"])
layers = [
    ("KNOWLEDGE LAYER", "Offline", "Turns official documentation into a searchable, "
     "semantic knowledge base of every module.", ACCENT),
    ("INTELLIGENCE LAYER", "Runtime", "An autonomous agent retrieves evidence, generates "
     "the playbook, and validates it.", ACCENT2),
    ("EXPERIENCE LAYER", "Interface", "A real-time chat that streams reasoning, shows "
     "sources, and keeps history.", OKGREEN),
]
cw = Inches(3.7); gap = Inches(0.3); y = Inches(2.65)
pos = []
for i, (h, tag, b, c) in enumerate(layers):
    x = Inches(0.8) + i * (cw + gap)
    rect(s, x, y, cw, Inches(3.4), color=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, y, cw, Inches(0.7), color=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y + Inches(0.35), cw, Inches(0.35), color=c)
    text(s, x + Inches(0.3), y + Inches(0.13), cw - Inches(0.5), Inches(0.45),
         [[R(h, 15, INK, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x + Inches(0.3), y + Inches(0.9), cw - Inches(0.6), Inches(0.4),
         [[R(tag.upper(), 12, c, True)]], space_after=0)
    text(s, x + Inches(0.3), y + Inches(1.4), cw - Inches(0.6), Inches(1.8),
         [[R(b, 13.5, SOFT)]], line_spacing=1.25, space_after=0)
    pos.append((x, cw))
for i in range(2):
    x0, w0 = pos[i]; x1, _ = pos[i + 1]
    connect(s, x0 + w0, y + Inches(1.7), x1, y + Inches(1.7), color=ACCENT)
text(s, Inches(0.8), Inches(6.45), Inches(11.6), Inches(0.5),
     [[R("Evidence flows left to right - the interface never shows an answer that "
         "skipped retrieval or validation.", 13.5, MUTED)]], space_after=0)
add_transition(s, "morph")

# =================================================================
# 8 - KNOWLEDGE FOUNDATION (the moat)
# =================================================================
s = slide()
set_bg(s)
kicker(s, "Our Moat", 8)
title(s, ["The advantage a general model cannot copy:",
          "a curated, refreshable knowledge base"])
rect(s, Inches(0.8), Inches(2.4), Inches(4.6), Inches(4.05), color=PANEL2,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
text(s, Inches(0.8), Inches(2.95), Inches(4.6), Inches(1.3),
     [[R("~1,230", 58, ACCENT, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
text(s, Inches(1.1), Inches(4.2), Inches(4.0), Inches(0.5),
     [[R("official modules indexed", 15, SOFT)]], align=PP_ALIGN.CENTER, space_after=0)
text(s, Inches(1.1), Inches(4.95), Inches(4.0), Inches(1.3),
     [[R("Five ecosystems:", 13, ACCENT2, True)],
      [R("Core  -  AWS  -  Azure  -  Kubernetes  -  Community", 13, SOFT)]],
     align=PP_ALIGN.CENTER, line_spacing=1.2, space_after=4)
feats = [
    ("Built from official sources", "Parsed into structured parameters, examples and return values."),
    ("Always refreshable", "Detects upstream changes and re-indexes - no training cutoff."),
    ("Semantic + auditable", "Stored as meaning, with a traceable link back to each source."),
]
for i, (h, b) in enumerate(feats):
    y = Inches(2.4) + i * Inches(1.38)
    rect(s, Inches(5.65), y, Inches(6.85), Inches(1.22), color=PANEL,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, Inches(5.65), y, Inches(0.09), Inches(1.22), color=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(5.95), y + Inches(0.18), Inches(6.3), Inches(0.45),
         [[R(h, 16, WHITE, True)]], space_after=0)
    text(s, Inches(5.95), y + Inches(0.64), Inches(6.35), Inches(0.5),
         [[R(b, 13, MUTED)]], line_spacing=1.05, space_after=0)
add_transition(s, "morph")

# =================================================================
# 9 - RETRIEVAL FUNNEL (diagram)
# =================================================================
s = slide()
set_bg(s)
kicker(s, "Grounding Engine", 9)
title(s, ["Retrieval is a funnel that finds the right evidence,",
          "not just text that looks similar"])
stages = [
    ("Analyze", "intent + explicit module"),
    ("Route", "to the right ecosystem(s)"),
    ("Hybrid search", "meaning + keywords"),
    ("Rerank", "by what the user wants"),
    ("Diversify", "no single source dominates"),
    ("Select", "top evidence to the agent"),
]
# funnel as narrowing stacked bars
widths = [Inches(11.7), Inches(10.6), Inches(9.5), Inches(8.4), Inches(7.3), Inches(6.2)]
y = Inches(2.4)
for i, (h, b) in enumerate(stages):
    w = widths[i]
    x = Inches(0.8) + (Inches(11.7) - w) / 2
    c = ACCENT if i < 2 else (ACCENT2 if i < 4 else OKGREEN)
    rect(s, x, y, w, Inches(0.62), color=PANEL, line=c, line_w=Pt(1.25),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x + Inches(0.12), y + Inches(0.13), Inches(0.36), Inches(0.36),
         color=c, shape=MSO_SHAPE.OVAL)
    text(s, x + Inches(0.12), y + Inches(0.11), Inches(0.36), Inches(0.36),
         [[R(str(i + 1), 13, INK, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x + Inches(0.6), y, Inches(3.2), Inches(0.62),
         [[R(h, 14.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x + Inches(3.5), y, w - Inches(3.7), Inches(0.62),
         [[R(b, 12, MUTED)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    if i < len(stages) - 1:
        cx = Inches(0.8) + Inches(11.7) / 2
        connect(s, cx, y + Inches(0.62), cx, y + Inches(0.74), color=MUTED, width=Pt(1.4))
    y = y + Inches(0.74)
text(s, Inches(0.8), Inches(6.95), Inches(11.6), Inches(0.4),
     [[R("Each stage narrows hundreds of candidates down to the few documents that truly matter.",
         12.5, MUTED, italic=True)]], align=PP_ALIGN.CENTER, space_after=0)
add_transition(s, "morph")

# =================================================================
# 10 - AGENT WORKFLOW (diagram, agent only)
# =================================================================
s = slide()
set_bg(s)
kicker(s, "The Autonomous Agent", 10)
title(s, ["The agent plans, acts with tools, and only then",
          "composes a grounded, verified reply"])
phases = [
    ("Plan", "read intent,\ndraft a tool plan", ACCENT),
    ("Execute", "call tools,\ngather evidence", ACCENT),
    ("Clarify", "ask only if\ndetails missing", GOLD),
    ("Generate", "write playbook\nfrom evidence", ACCENT),
    ("Synthesize", "compose grounded\nfinal answer", OKGREEN),
]
cw = Inches(2.3); gap = Inches(0.16); y = Inches(2.5)
pos = []
for i, (h, b, c) in enumerate(phases):
    x = Inches(0.7) + i * (cw + gap)
    rect(s, x, y, cw, Inches(1.85), color=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, y, cw, Inches(0.1), color=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, y + Inches(0.24), cw, Inches(0.4),
         [[R("0%d" % (i + 1), 15, c, True)]], align=PP_ALIGN.CENTER, space_after=0)
    text(s, x, y + Inches(0.62), cw, Inches(0.4),
         [[R(h, 15.5, WHITE, True)]], align=PP_ALIGN.CENTER, space_after=0)
    text(s, x + Inches(0.1), y + Inches(1.05), cw - Inches(0.2), Inches(0.7),
         [[R(b, 11, MUTED)]], align=PP_ALIGN.CENTER, line_spacing=1.0, space_after=0)
    pos.append((x, cw))
for i in range(len(pos) - 1):
    x0, w0 = pos[i]; x1, _ = pos[i + 1]
    connect(s, x0 + w0, y + Inches(0.9), x1, y + Inches(0.9), color=MUTED, width=Pt(1.5))
text(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(0.4),
     [[R("Clarify is conditional - skipped when the request is already complete.",
         12, MUTED, italic=True)]], align=PP_ALIGN.CENTER, space_after=0)
# tools + intents
by = Inches(5.05); bw = Inches(5.78)
rect(s, Inches(0.7), by, bw, Inches(1.75), color=PANEL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
text(s, Inches(0.95), by + Inches(0.16), bw - Inches(0.5), Inches(0.4),
     [[R("Tools it can call", 14, ACCENT, True)]], space_after=0)
tools = ["Search the documentation", "Inspect a module's reference",
         "Validate a YAML playbook", "Generate a playbook"]
for i, t in enumerate(tools):
    col = i % 2; row = i // 2
    tx = Inches(0.95) + col * Inches(2.78)
    ty = by + Inches(0.64) + row * Inches(0.5)
    rect(s, tx, ty + Inches(0.06), Inches(0.13), Inches(0.13), color=ACCENT2, shape=MSO_SHAPE.OVAL)
    text(s, tx + Inches(0.27), ty, Inches(2.55), Inches(0.4),
         [[R(t, 12, SOFT)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
ix = Inches(6.68)
rect(s, ix, by, bw, Inches(1.75), color=PANEL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
text(s, ix + Inches(0.25), by + Inches(0.16), bw - Inches(0.5), Inches(0.4),
     [[R("Intents it understands", 14, ACCENT, True)]], space_after=0)
intents = ["Generate", "Explain", "Troubleshoot", "Compare", "Edit", "Chat"]
for i, t in enumerate(intents):
    col = i % 3; row = i // 3
    px = ix + Inches(0.25) + col * Inches(1.8)
    py = by + Inches(0.68) + row * Inches(0.5)
    chip(s, px, py, Inches(1.62), t, color=ACCENT2, h=Inches(0.4), size=11.5)
add_transition(s, "morph")

# =================================================================
# 11 - TRUST LAYER (validation)
# =================================================================
s = slide()
set_bg(s)
kicker(s, "Trust by Construction", 11)
title(s, ["The safety gate a general chatbot simply does not have"])
text(s, Inches(0.8), Inches(2.05), Inches(11.6), Inches(0.5),
     [[R("Before any playbook is shown, it must pass every check below. "
         "If one fails, the agent revises and tries again.", 14.5, SOFT)]], space_after=0)
checks = [
    "Valid YAML syntax", "Sound playbook structure", "Platform layout rules",
    "Every module truly exists", "Required parameters present", "Targets are defined",
    "No leftover placeholders", "No hard-coded secrets", "Best-practice linting",
]
cw = Inches(3.78); ch = Inches(0.92); gx = Inches(0.25); gy = Inches(0.22)
for i, t in enumerate(checks):
    col = i % 3; row = i // 3
    x = Inches(0.8) + col * (cw + gx)
    y = Inches(2.75) + row * (ch + gy)
    rect(s, x, y, cw, ch, color=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x + Inches(0.22), y + Inches(0.26), Inches(0.42), Inches(0.42),
         color=OKGREEN, shape=MSO_SHAPE.OVAL)
    text(s, x + Inches(0.22), y + Inches(0.24), Inches(0.42), Inches(0.42),
         [[R("OK", 10.5, INK, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x + Inches(0.82), y, cw - Inches(1.0), ch,
         [[R(t, 14, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
add_transition(s, "morph")

# =================================================================
# 12 - THE EXPERIENCE (UI placeholders)
# =================================================================
s = slide()
set_bg(s)
kicker(s, "The Experience", 12)
title(s, ["The reasoning is visible - so the user can trust what they see"])
img_placeholder(s, Inches(0.8), Inches(2.35), Inches(5.75), Inches(3.95),
                "Screenshot: the chat view generating a playbook - show the live step "
                "progress (Plan / Search / Write / Validate) and the YAML result card.",
                art="Screenshot")
img_placeholder(s, Inches(6.75), Inches(2.35), Inches(5.75), Inches(3.95),
                "Screenshot: the source attribution and validation report - the cited "
                "modules and the green pass / warning breakdown.",
                art="Screenshot")
text(s, Inches(0.8), Inches(6.45), Inches(11.6), Inches(0.6),
     [[R("Live progress  -  cited sources  -  validation report  -  threaded history  -  usage analytics",
         14, SOFT, True)]], align=PP_ALIGN.CENTER, space_after=0)
add_transition(s, "morph")

# =================================================================
# 13 - COMPETITIVE MATRIX
# =================================================================
s = slide()
set_bg(s)
kicker(s, "Why Choose This", 13)
title(s, ["Side by side with the alternatives"])
cols = ["Capability", "Manual\nexpert", "ChatGPT /\nClaude", "This\nproject"]
rows = [
    ("Plain-English to playbook", "No", "Yes", "Yes"),
    ("Grounded in official docs", "Yes", "No", "Yes"),
    ("Current, refreshable knowledge", "Yes", "No", "Yes"),
    ("Automatic validation gate", "Manual", "No", "Yes"),
    ("Cites its sources", "n/a", "No", "Yes"),
    ("Resists hallucinated modules", "Yes", "No", "Yes"),
    ("Fast & accessible to non-experts", "No", "Partly", "Yes"),
]
tx = Inches(0.8); ty = Inches(2.2)
c0 = Inches(5.0); cN = Inches(2.18); rh = Inches(0.56); hh = Inches(0.7)
colx = [tx, tx + c0, tx + c0 + cN, tx + c0 + 2 * cN]
colw = [c0, cN, cN, cN]
# header
for j, ctitle in enumerate(cols):
    c = ACCENT if j == 3 else (PANEL2 if j == 0 else PANEL)
    hl = INK if j == 3 else (ACCENT if j == 0 else SOFT)
    rect(s, colx[j], ty, colw[j], hh, color=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, colx[j], ty, colw[j], hh, [[R(cols[j], 13, hl, True)]],
         align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER), anchor=MSO_ANCHOR.MIDDLE,
         space_after=0, line_spacing=0.95)
    if j == 0:
        # re-pad label
        pass
def verdict_color(v):
    if v == "Yes":
        return OKGREEN
    if v == "No":
        return ERRRED
    return MUTED
for i, row in enumerate(rows):
    y = ty + hh + Inches(0.08) + i * (rh + Inches(0.06))
    rect(s, colx[0], y, c0, rh, color=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, colx[0] + Inches(0.25), y, c0 - Inches(0.4), rh,
         [[R(row[0], 13.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    for j in range(1, 4):
        v = row[j]
        own = (j == 3)
        bg = PANEL2 if own else PANEL
        rect(s, colx[j], y, colw[j], rh, color=bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        if own:
            rect(s, colx[j], y, Inches(0.06), rh, color=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, colx[j], y, colw[j], rh, [[R(v, 13, verdict_color(v), True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
add_transition(s, "morph")

# =================================================================
# 14 - PROOF / RESULTS
# =================================================================
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

s = slide()
set_bg(s)
kicker(s, "Proof", 14)
title(s, ["Measured on a transparent benchmark - and honest about it"])
big_stat(s, Inches(0.8), Inches(2.35), Inches(3.0), "78", "average score / 100", ACCENT)
big_stat(s, Inches(4.0), Inches(2.35), Inches(3.0), "60%", "of cases pass the bar", OKGREEN)
big_stat(s, Inches(7.2), Inches(2.35), Inches(3.0), "30", "golden test scenarios", ACCENT2)
rect(s, Inches(0.8), Inches(4.65), Inches(9.4), Inches(2.5), color=PANEL,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
text(s, Inches(1.05), Inches(4.8), Inches(9), Inches(0.4),
     [[R("Average performance across five evaluation layers", 13, SOFT, True)]], space_after=0)
cd = CategoryChartData()
cd.categories = ["Intent", "Retrieval", "Module", "Playbook", "Runtime"]
cd.add_series("Score", (88, 64, 82, 80, 90))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(1.0), Inches(5.2), Inches(9.0), Inches(1.9), cd)
ch = gf.chart
ch.has_legend = False; ch.has_title = False
ch.plots[0].gap_width = 80
ser = ch.plots[0].series[0]
ser.format.fill.solid(); ser.format.fill.fore_color.rgb = ACCENT
ca = ch.category_axis
ca.tick_labels.font.size = Pt(11); ca.tick_labels.font.color.rgb = SOFT
ca.format.line.color.rgb = PANEL2
va = ch.value_axis
va.maximum_scale = 100; va.minimum_scale = 0; va.has_major_gridlines = False
va.tick_labels.font.size = Pt(10); va.tick_labels.font.color.rgb = MUTED
va.format.line.fill.background()
rect(s, Inches(10.45), Inches(2.35), Inches(2.05), Inches(4.8), color=PANEL2,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
text(s, Inches(10.65), Inches(2.55), Inches(1.7), Inches(4.5),
     [[R("Reading", 14, ACCENT, True)], [R("", 6, SOFT)],
      [R("Strong on mainstream cloud and core tasks.", 12.5, SOFT)], [R("", 5, SOFT)],
      [R("Retrieval on the hardest ecosystem is the clearest lever to improve.", 12.5, SOFT)]],
     line_spacing=1.15, space_after=2)
add_transition(s, "morph")

# =================================================================
# 15 - IMPACT / WHO IT IS FOR
# =================================================================
s = slide()
set_bg(s)
kicker(s, "Impact", 15)
title(s, ["Who benefits - and the value it unlocks"])
audiences = [
    ("DevOps & SRE teams", "Ship automation faster, with fewer runtime failures and a built-in safety net.", ACCENT),
    ("Engineers learning IaC", "A guided, explainable on-ramp - every answer cites the docs to learn from.", ACCENT2),
    ("Organizations", "Lower the expertise barrier and standardize on validated, auditable output.", OKGREEN),
]
cw = Inches(3.8); gap = Inches(0.25); y = Inches(2.5)
for i, (h, b, c) in enumerate(audiences):
    x = Inches(0.8) + i * (cw + gap)
    rect(s, x, y, cw, Inches(2.3), color=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, x, y, cw, Inches(0.1), color=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + Inches(0.3), y + Inches(0.3), cw - Inches(0.6), Inches(0.5),
         [[R(h, 16.5, WHITE, True)]], space_after=0)
    text(s, x + Inches(0.3), y + Inches(0.9), cw - Inches(0.6), Inches(1.3),
         [[R(b, 13.5, SOFT)]], line_spacing=1.25, space_after=0)
vals = [("Faster", "minutes, not hours"), ("Safer", "verified before use"),
        ("Explainable", "sources on every answer"), ("Accessible", "no deep expertise needed")]
for i, (h, b) in enumerate(vals):
    x = Inches(0.8) + i * Inches(3.02)
    yy = Inches(5.25)
    rect(s, x, yy, Inches(2.85), Inches(1.3), color=PANEL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    text(s, x + Inches(0.25), yy + Inches(0.22), Inches(2.4), Inches(0.45),
         [[R(h, 16, ACCENT, True)]], space_after=0)
    text(s, x + Inches(0.25), yy + Inches(0.7), Inches(2.45), Inches(0.5),
         [[R(b, 12, MUTED)]], space_after=0)
add_transition(s, "morph")

# =================================================================
# 16 - LIMITATIONS & ROADMAP
# =================================================================
s = slide()
set_bg(s)
kicker(s, "Honest Limits & Roadmap", 16)
title(s, ["What is not solved yet - and where it goes next"])
lx = Inches(0.8); ly = Inches(2.3); lw = Inches(5.6)
rect(s, lx, ly, lw, Inches(4.3), color=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
rect(s, lx, ly, lw, Inches(0.08), color=GOLD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, lx + Inches(0.3), ly + Inches(0.25), lw - Inches(0.6), Inches(0.5),
     [[R("Current limitations", 17, GOLD, True)]], space_after=0)
limits = ["Retrieval is weaker on the most complex ecosystem",
          "Generation is non-deterministic by nature",
          "Linting depends on a compatible environment",
          "No automated delivery pipeline yet"]
for i, t in enumerate(limits):
    text(s, lx + Inches(0.3), ly + Inches(0.95) + i * Inches(0.78), lw - Inches(0.6), Inches(0.7),
         [[R("-  ", 14, GOLD, True), R(t, 13.5, SOFT)]], line_spacing=1.1, space_after=0)
px = Inches(6.75)
rect(s, px, ly, lw, Inches(4.3), color=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
rect(s, px, ly, lw, Inches(0.08), color=OKGREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, px + Inches(0.3), ly + Inches(0.25), lw - Inches(0.6), Inches(0.5),
     [[R("Where it goes next", 17, OKGREEN, True)]], space_after=0)
persp = ["Strengthen retrieval on hard ecosystems",
         "Containerize and add continuous delivery",
         "Broaden coverage to more ecosystems",
         "Explore fine-tuning and multi-step scenarios"]
for i, t in enumerate(persp):
    text(s, px + Inches(0.3), ly + Inches(0.95) + i * Inches(0.78), lw - Inches(0.6), Inches(0.7),
         [[R("+  ", 14, OKGREEN, True), R(t, 13.5, SOFT)]], line_spacing=1.1, space_after=0)
add_transition(s, "morph")

# =================================================================
# 17 - CONCLUSION
# =================================================================
s = slide()
set_bg(s, BG1, BG2, angle=60)
rect(s, Inches(8.8), Inches(2.5), Inches(6.2), Inches(6.2), color=ACCENT, alpha=7, shape=MSO_SHAPE.OVAL)
rect(s, Inches(0.92), Inches(0.95), Inches(0.9), Inches(0.10), color=ACCENT)
text(s, Inches(0.9), Inches(1.2), Inches(10), Inches(0.5),
     [[R("CONCLUSION", 14, ACCENT, True)]], space_after=0)
text(s, Inches(0.85), Inches(1.75), Inches(11.6), Inches(1.6),
     [[R("Generic AI is fluent. This is ", 34, WHITE, True), R("trustworthy", 34, ACCENT, True),
       R(".", 34, WHITE, True)],
      [R("Grounded, verified, and explainable IaC - in plain English.", 24, SOFT, False)]],
     line_spacing=1.1, space_after=4)
takeaways = [
    "It retrieves real documentation instead of guessing from memory.",
    "It verifies every playbook before you ever run it.",
    "It explains itself - every answer is backed by cited sources.",
]
for i, t in enumerate(takeaways):
    y = Inches(3.85) + i * Inches(0.62)
    rect(s, Inches(0.95), y + Inches(0.08), Inches(0.16), Inches(0.16), color=ACCENT2, shape=MSO_SHAPE.OVAL)
    text(s, Inches(1.3), y, Inches(10.8), Inches(0.55),
         [[R(t, 15.5, SOFT)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
rect(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(0.02), color=PANEL2)
text(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.6),
     [[R("Thank you - questions and feedback welcome.", 18, WHITE, True)]], space_after=0)
text(s, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.4),
     [[R("Final-Year Engineering Project (PFE)  -  2026", 13, MUTED)]], space_after=0)
add_transition(s, "fade")

# =================================================================
# SPEAKER NOTES
# =================================================================
SLIDE_NOTES = {
    1: ("Open with confidence and the core claim: generic AI guesses at infrastructure; this project "
        "proves it. One line: an assistant that writes production automation from plain English, grounded "
        "in official docs and verified before you run it. Introduce yourself and the project context."),
    2: ("Establish the stakes so the audience feels the problem. Infrastructure as Code runs the cloud, "
        "and it is unforgiving: a wrong module means a failed deploy, a wrong parameter can open a security "
        "hole, and either way engineers lose hours debugging. This is a high-stakes domain where guessing "
        "is expensive - which sets up why a general chatbot is risky here."),
    3: ("State the problematic as a single, sharp research question and pause on it. Then break it into the "
        "three things that must hold at once: correctness (does it run?), trust (can we audit it?), and "
        "accessibility (can a non-expert use it safely?). The rest of the talk shows how we satisfy all three."),
    4: ("This is the competitive crux - take your time. ChatGPT and Claude are fluent, but they answer from "
        "memory. In this domain that gives six concrete failure modes: hallucinated modules and parameters, "
        "no grounding in the real docs, a training-knowledge cutoff, no validation so errors surface at "
        "runtime, no provenance to audit, and no specialization. The issue isn't fluency - it's the absence "
        "of evidence and verification."),
    5: ("Pivot from problem to promise. Our answer turns a language model into a trustworthy domain expert "
        "through three pillars: GROUNDED (retrieves official docs first), VERIFIED (a validation gate catches "
        "errors before delivery), and TRANSPARENT (cites the exact modules used). These are precisely the "
        "guarantees a general chatbot cannot make."),
    6: ("Use the diagram to make the difference visceral. A general chatbot goes request to LLM-from-memory to "
        "a plausible answer you only test at runtime. Our path inserts the two steps that matter: retrieve the "
        "official documentation, generate from that evidence, then validate - delivering a trusted result with "
        "cited sources. The 'detour' through docs and validation is the entire value."),
    7: ("Zoom out to the system. Three layers: a knowledge layer that builds the searchable doc base offline; "
        "an intelligence layer where the autonomous agent retrieves, generates and validates; and an experience "
        "layer - the real-time chat. Stress that evidence flows left to right and the UI never shows an answer "
        "that skipped retrieval or validation."),
    8: ("Frame the knowledge base as the moat - the thing a general model cannot replicate. About 1,230 official "
        "modules across five ecosystems, built from official sources, refreshable so there is no training cutoff, "
        "and stored semantically with a traceable link back to each source. This is what makes grounding possible."),
    9: ("Explain grounding as a funnel that narrows hundreds of candidates to the few documents that matter: "
        "analyze the query, route to the right ecosystem, hybrid search (meaning plus keywords), rerank by intent, "
        "diversify so no single source dominates, and select the top evidence. Quality retrieval is what keeps the "
        "generated code honest."),
    10: ("Walk the autonomous agent loop using the diagram. Plan: read intent and draft a tool plan. Execute: call "
         "tools and gather evidence - the ecosystem to search is chosen by a retrieval vote over the embeddings, not "
         "by trusting the model's guess. Clarify: only if essential details are missing. Generate: write the playbook "
         "from the gathered evidence and the user's parameters. Synthesize: compose the grounded final reply. It "
         "exposes four tools and understands six intents."),
    11: ("This is the gate a chatbot lacks - emphasize it. Before any playbook is shown it must pass every check: "
         "valid YAML, sound structure, platform rules, the module actually exists, required parameters present, "
         "targets defined, no placeholders, no hard-coded secrets, and best-practice linting. On failure the agent "
         "revises and retries. Trust is built in, not hoped for."),
    12: ("Show the product. Replace the two placeholders with real screenshots: one of the chat generating a playbook "
         "with live step progress and the YAML card, one of the cited sources and the validation report. The point: "
         "the reasoning is visible, which is what earns user trust."),
    13: ("Land the decision with the matrix. Read down the 'This project' column - it is the only one that is Yes on "
         "grounding, current knowledge, validation, citations and resisting hallucination, while still being fast and "
         "accessible. Manual experts are safe but slow; general AI is fast but unverified. We are both safe and fast."),
    14: ("Be rigorous and honest - this is an academic defense. Thirty curated scenarios, scored on five weighted "
         "layers, passing at 70/100. Average 78, 60% pass. The chart shows strength on intent, modules, playbook "
         "quality and runtime; retrieval on the hardest ecosystem is the weakest layer and the clearest place to "
         "improve. Owning the weakness builds credibility."),
    15: ("Broaden to value and audience. DevOps and SRE teams ship faster with fewer failures; engineers learning IaC "
         "get an explainable on-ramp; organizations lower the expertise barrier and standardize on validated, auditable "
         "output. Summarize the value: faster, safer, explainable, accessible."),
    16: ("Stay credible: name real limits - retrieval on the hardest ecosystem, non-deterministic generation, linting "
         "environment dependence, and no delivery pipeline yet - then the roadmap: strengthen retrieval, containerize "
         "and add CI/CD, broaden coverage, and explore fine-tuning and multi-step scenarios."),
    17: ("Close on the thesis: generic AI is fluent, this is trustworthy - grounded, verified, explainable IaC in plain "
         "English. Recap the three proofs, then open the floor for questions. Keep this slide on screen during Q&A."),
}
for _i, _sl in enumerate(prs.slides, start=1):
    if _i in SLIDE_NOTES:
        notes(_sl, SLIDE_NOTES[_i])

# =================================================================
import os

out = r"c:\Users\ahmed\OneDrive\Desktop\pfe2026\ansible-iac-ai\AnsibleAI_Presentation.pptx"
try:
    prs.save(out)
except PermissionError:
    out = out.replace(".pptx", "_v2.pptx")
    prs.save(out)
    print("Main file was locked (open in PowerPoint). Saved alternate copy.")
# clean up the stale alternate if we saved the main file successfully
alt = r"c:\Users\ahmed\OneDrive\Desktop\pfe2026\ansible-iac-ai\AnsibleAI_Presentation_v2.pptx"
if out.endswith("AnsibleAI_Presentation.pptx") and os.path.exists(alt):
    try:
        os.remove(alt)
    except OSError:
        pass
print("SAVED:", out, "slides:", len(prs.slides._sldIdLst))
